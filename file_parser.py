import os
import zipfile
import tempfile
import shutil
from pathlib import Path

def parse_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.pdf':
            return parse_pdf(file_path)
        elif ext == '.docx':
            return parse_docx(file_path)
        elif ext == '.xlsx':
            return parse_xlsx(file_path)
        elif ext == '.pptx':
            return parse_pptx(file_path)
        elif ext == '.odt':
            return parse_odt(file_path)
        elif ext in ['.txt', '.py', '.js', '.json', '.md', '.csv', '.xml', '.html', '.css', '.log']:
            return parse_text(file_path)
        elif ext == '.zip':
            return parse_zip(file_path)
        elif ext == '.rar':
            return parse_rar(file_path)
        else:
            return f"Файл формата {ext} не поддерживается для автоматического чтения."
    except Exception as e:
        return f"Ошибка при парсинге файла: {e}"

def parse_pdf(file_path):
    import PyPDF2
    reader = PyPDF2.PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text.strip() or "PDF не содержит текста."

def parse_docx(file_path):
    from docx import Document
    doc = Document(file_path)
    full_text = [para.text for para in doc.paragraphs]
    return "\n".join(full_text).strip() or "DOCX пуст."

def parse_xlsx(file_path):
    from openpyxl import load_workbook
    wb = load_workbook(file_path, data_only=True)
    all_text = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                rows.append(" ".join(row_text))
        if rows:
            all_text.append(f"Лист: {sheet.title}\n" + "\n".join(rows))
    return "\n\n".join(all_text).strip() or "XLSX пуст."

def parse_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs).strip() or "PPTX пуст."

def parse_odt(file_path):
    from odf import text, teletype
    from odf.opendocument import load
    doc = load(file_path)
    all_paras = doc.getElementsByType(text.P)
    content = [teletype.extractText(p) for p in all_paras]
    return "\n".join(content).strip() or "ODT пуст."

def parse_text(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip() or "Файл пуст."

def parse_zip(file_path):
    temp_dir = tempfile.mkdtemp()
    content = []
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            z.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for f in files:
                full_path = os.path.join(root, f)
                try:
                    with open(full_path, "r", encoding="utf-8", errors="ignore") as fp:
                        data = fp.read()
                        if data.strip():
                            content.append(f"Файл: {f}\n{data[:1000]}" + ("..." if len(data) > 1000 else ""))
                except:
                    content.append(f"Не удалось прочитать файл: {f} (возможно, бинарный)")
        return "\n\n".join(content) if content else "Архив не содержит читаемых текстовых файлов."
    except Exception as e:
        return f"Ошибка при распаковке архива: {e}"
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def parse_rar(file_path):
    try:
        import rarfile
    except ImportError:
        return "Поддержка RAR не установлена. Установите библиотеку rarfile и утилиту unrar."
    temp_dir = tempfile.mkdtemp()
    content = []
    try:
        with rarfile.RarFile(file_path, 'r') as rf:
            rf.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for f in files:
                full_path = os.path.join(root, f)
                try:
                    with open(full_path, "r", encoding="utf-8", errors="ignore") as fp:
                        data = fp.read()
                        if data.strip():
                            content.append(f"Файл: {f}\n{data[:1000]}" + ("..." if len(data) > 1000 else ""))
                except:
                    content.append(f"Не удалось прочитать файл: {f} (возможно, бинарный)")
        return "\n\n".join(content) if content else "Архив не содержит читаемых текстовых файлов."
    except Exception as e:
        return f"Ошибка при распаковке RAR: {e}"
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)